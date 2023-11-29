"""parsing text"""
import io
import os
import re
import csv
import fitz
import docx
import json
import ezodf
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text


class TextExtractor:
    def __init__(self, file_path):
        self.file_path = file_path

    def __enter__(self):
        _, self.extension = os.path.splitext(self.file_path)
        self.extension = self.extension.lower()

        try:
            if self.extension != ".rtf":
                with open(self.file_path, 'rb') as file:
                    self.file_content = file.read()
            else:
                with open(self.file_path, 'r') as file:
                    self.file_content = file.read()
        except Exception as e:
            raise RuntimeError(f"Error reading file: {str(e)}")

        return self

    def __exit__(self, exc_type, exc_value, traceback):
        pass

    def extract_text(self):
        if self.extension == '.pdf':
            return self.extract_text_from_pdf()
        elif self.extension == '.docx':
            return self.extract_text_from_word()
        elif self.extension == '.pptx':
            return self.extract_text_from_powerpoint()
        elif self.extension == '.rtf':
            return self.extract_text_from_rtf()
        elif self.extension == '.xml':
            return self.extract_text_from_xml()
        elif self.extension == '.odt':
            return self.extract_text_from_odt()
        elif self.extension == '.csv':
            return self.extract_text_from_csv()
        elif self.extension == '.xlsx':
            return self.extract_text_from_xlsx()
        elif self.extension == '.json':
            return self.extract_text_from_json()
        elif self.extension in ('.txt', '.log'):
            return self.extract_text_from_other()
        else:
            return "Unsupported file format"

    def extract_text_from_pdf(self):
        try:
            doc = fitz.Document(self.file_path)
            text = ''
            for page_num in range(doc.page_count):
                page = doc[page_num]
                text += page.get_text()
            return text
        except Exception as e:
            return f"Error extracting text from PDF: {str(e)}"

    def extract_text_from_word(self):
        try:
            doc = docx.Document(io.BytesIO(self.file_content))
            text = ''
            for paragraph in doc.paragraphs:
                text += paragraph.text + '\n'
            return text
        except Exception as e:
            return f"Error extracting text from Word document: {str(e)}"

    def extract_text_from_rtf(self):
        try:
            text = rtf_to_text(self.file_content)
            return text
        except Exception as e:
            return f"Error extracting text from RTF document: {str(e)}"

    def extract_text_from_powerpoint(self):
        try:
            presentation = Presentation(io.BytesIO(self.file_content))
            text = ''
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
            return text
        except Exception as e:
            return f"Error extracting text from PowerPoint presentation: {str(e)}"

    def extract_text_from_xml(self):
        try:
            with open(self.file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'xml')

            def extract_tags_info(tag, depth=0):
                tag_info = f"{'  ' * depth}{tag.name}: {tag.get_text(strip=True)}"
                for child in tag.find_all(recursive=False):
                    tag_info += f"\n{extract_tags_info(child, depth + 1)}"
                return tag_info

            result_str = extract_tags_info(soup)
            return result_str
        except Exception as e:
            return f"Error extracting text from XML: {str(e)}"

    def extract_text_from_odt(self):
        try:
            odt = ezodf.opendoc(self.file_path)
            list = ""
            # Запускаем цикл for  и перебираем все что нашли в файле)
            for i in odt.body:
                if i.text == None:
                    continue
                else:
                    list += " ".join(re.findall(r"[\w']+", i.text.lower()))
            return list.strip()
        except Exception as e:
            return f"Error extracting text from ODT: {str(e)}"

    def extract_text_from_csv(self):
        try:
            with open(self.file_path, 'r', encoding='utf-8') as file:
                reader = csv.reader(file)
                text = ' '.join([' '.join(row) for row in reader])
            return text.strip()
        except Exception as e:
            return f"Error extracting text from CSV: {str(e)}"

    def extract_text_from_xlsx(self):
        try:
            workbook = openpyxl.load_workbook(self.file_path)
            # Создание словаря для хранения данных из всех листов
            all_data = {}

            # Чтение данных из каждого листа и добавление их в словарь
            for sheet_name in workbook.sheetnames:
                # Выбор текущего листа
                sheet = workbook[sheet_name]

                # Получение данных из листа
                data = []
                for row in sheet.iter_rows(values_only=True):
                    data.append(row)

                # Добавление данных в словарь
                all_data[sheet_name] = data
            return "\n".join((f"{key}: {value}" for key, value in all_data.items()))
        except Exception as e:
            return f"Error extracting text from XLSX: {str(e)}"

    def extract_text_from_json(self):
        try:
            with open(self.file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            text = json.dumps(data, indent=2)
            return text.strip()
        except Exception as e:
            return f"Error extracting text from JSON: {str(e)}"

    def extract_text_from_other(self):
        # For other file formats, you can add specific parsing logic here.
        # This is just a placeholder for plain text files.
        return self.file_content.decode('utf-8', errors='ignore')
